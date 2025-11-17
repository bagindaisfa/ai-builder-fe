import os
import uuid
import shutil
import math
from datetime import datetime
from pathlib import Path
from typing import List, Dict, Optional, BinaryIO
import mimetypes
from werkzeug.datastructures import FileStorage
from app import db
from app.models import Document
from app.models.document_vector import DocumentVector

class DocumentService:
    """
    Service for handling document storage and retrieval operations.
    Handles file system operations for document management.
    """
    
    def __init__(self, base_storage_path: str = None):
        """
        Initialize the DocumentService with a base storage path.
        
        Args:
            base_storage_path: Base directory where documents will be stored.
                              If not provided, uses environment variable DOCUMENT_STORAGE_PATH
                              or defaults to 'data/documents'.
        """
        self.base_storage_path = (
            base_storage_path or 
            os.environ.get('DOCUMENT_STORAGE_PATH') or 
            os.path.join(os.getcwd(), 'data', 'documents')
        )
        
        # Ensure base storage directory exists
        os.makedirs(self.base_storage_path, exist_ok=True)
    
    def generate_unique_filename(self, original_filename: str) -> str:
        """
        Generate a unique filename while preserving the original extension.
        
        Args:
            original_filename: The original filename
            
        Returns:
            str: A unique filename with preserved extension
        """
        ext = os.path.splitext(original_filename)[1].lower()
        return f"{uuid.uuid4().hex}{ext}"
    
    def get_storage_path(self, subdirectory: str = '') -> str:
        """
        Get the full path to a storage subdirectory, creating it if it doesn't exist.
        
        Args:
            subdirectory: Subdirectory within the base storage path
            
        Returns:
            str: Full path to the storage directory
        """
        path = os.path.join(self.base_storage_path, subdirectory)
        os.makedirs(path, exist_ok=True)
        return path
    
    def save_document(
        self, 
        file: FileStorage, 
        subdirectory: str = '',
        custom_filename: str = None
    ) -> Dict[str, str]:
        """
        Save an uploaded file to the storage directory.
        
        Args:
            file: The file to save (Werkzeug FileStorage object)
            subdirectory: Subdirectory within the base storage path
            custom_filename: Optional custom filename (without path)
            
        Returns:
            Dict containing file metadata including:
            - id: Unique identifier for the file
            - original_filename: Original filename
            - stored_filename: Name under which the file was stored
            - file_path: Full path to the stored file
            - file_size: Size of the file in bytes
            - mime_type: Detected MIME type
            - created_at: Timestamp of when the file was saved
        """
        if not file or not file.filename:
            raise ValueError("No file provided")
        
        # Prepare file metadata
        original_filename = file.filename
        file_ext = os.path.splitext(original_filename)[1].lower()
        stored_filename = custom_filename or self.generate_unique_filename(original_filename)
        
        # Determine storage path
        storage_dir = self.get_storage_path(subdirectory)
        print(storage_dir)
        print(stored_filename)
        file_path = os.path.join(storage_dir, stored_filename)
        print(file_path)
        
        try:
            # Save the file
            file.save(file_path)
        except Exception as e:
            print (f"Failed to save file: {str(e)}")
            raise ValueError(f"Failed to save file: {str(e)}")
        
        # Get file size
        file_size = os.path.getsize(file_path)
        
        # Detect MIME type
        mime_type, _ = mimetypes.guess_type(original_filename)
        
        doc_info = {
            'id': str(uuid.uuid4()),
            'filename': original_filename,
            'stored_filename': stored_filename,
            'file_path': file_path,
            'file_size': file_size,
            'mime_type': mime_type or 'application/octet-stream',
            'created_at': datetime.utcnow().isoformat(),
            'file_extension': file_ext
        }

        try:
            # Open the file in binary mode for _extract_text
            with open(file_path, 'rb') as f:
                content_file = self._extract_text(f, mime_type or '')

            # Save document to database
            document = Document(
                uuid=doc_info['id'],
                filename=original_filename,
                content_type=mime_type,
                content=content_file,
                embedding_status='pending',
                path=file_path,
                stored_filename=stored_filename,
                size=file_size,
                mime_type=mime_type,
                extension=file_ext,
                doc_metadata=doc_info,
                word_count=len(content_file.split()) if content_file else 0,
                created_at=datetime.utcnow(),
                updated_at=datetime.utcnow()
            )
            db.session.add(document)
            db.session.commit()
        except Exception as e:
            db.session.rollback()
            # Clean up the saved file if database operation fails
            if os.path.exists(file_path):
                os.remove(file_path)
            raise e
        
        return doc_info
    
    def save_documents(
        self, 
        files: List[FileStorage], 
        subdirectory: str = ''
    ) -> List[Dict[str, str]]:
        """
        Save multiple uploaded files to the storage directory.
        
        Args:
            files: List of files to save (Werkzeug FileStorage objects)
            subdirectory: Subdirectory within the base storage path
            
        Returns:
            List of file metadata dictionaries
        """
        if not files:
            raise ValueError("No files provided")
            
        results = []
        
        for file in files:
            try:
                result = self.save_document(file, subdirectory)
                results.append({
                    'success': True,
                    'data': result
                })
            except Exception as e:
                results.append({
                    'success': False,
                    'error': str(e),
                    'filename': file.filename if file else 'unknown'
                })
        
        return results
    
    def get_document_path(self, filename: str, subdirectory: str = '') -> Optional[str]:
        """
        Get the full path to a stored document.
        
        Args:
            filename: Name of the file to retrieve
            subdirectory: Subdirectory within the base storage path
            
        Returns:
            Full path to the file if it exists, None otherwise
        """
        if not filename:
            return None
            
        file_path = os.path.join(self.get_storage_path(subdirectory), filename)
        return file_path if os.path.exists(file_path) else None


    def get_document_path_by_doc_id(self, uuid: str, subdirectory: str = '') -> Optional[str]:
        """
        Get the full path to a stored document.
        
        Args:
            filename: Name of the file to retrieve
            subdirectory: Subdirectory within the base storage path
            
        Returns:
            Full path to the file if it exists, None otherwise
        """
        document = Document.query.get(uuid)
        if not document:
            return None
            
        file_path = document.file_path
        return file_path if os.path.exists(file_path) else None
    
    def delete_document(self, filename: str, subdirectory: str = '') -> bool:
        """
        Delete a stored document.
        
        Args:
            filename: Name of the file to delete
            subdirectory: Subdirectory within the base storage path
            
        Returns:
            bool: True if deletion was successful, False otherwise
        """
        file_path = self.get_document_path(filename, subdirectory)
        if not file_path:
            return False
            
        try:
            os.remove(file_path)
            return True
        except OSError:
            return False

    
    def delete_document_by_uuid(self, uuid: str, subdirectory: str = '') -> bool:
        """
        Delete a stored document.
        
        Args:
            filename: Name of the file to delete
            subdirectory: Subdirectory within the base storage path
            
        Returns:
            bool: True if deletion was successful, False otherwise
        
        """

        ##document = Document.query.get(uuid)
        document = Document.query.filter_by(uuid=uuid).first()
        if not document:
            return False

        """
        get document_vector
        """
        document_vectors = DocumentVector.query.filter_by(document_uuid=uuid).all()
        if document_vectors:
            for document_vector in document_vectors:
                db.session.delete(document_vector)


        file_path = document.path
        if not file_path:
            return False
            
        """
        r4emove from db
        """
        try:

            db.session.delete(document)
            db.session.commit()
            os.remove(file_path)
            return True
        except OSError:
            return False
    
    def list_documents(self, subdirectory: str = '') -> List[Dict[str, str]]:
        """
        List all documents in a subdirectory.
        
        Args:
            subdirectory: Subdirectory within the base storage path
            
        Returns:
            List of document metadata dictionaries
        """
        dir_path = self.get_storage_path(subdirectory)
        documents = []
        
        for filename in os.listdir(dir_path):
            file_path = os.path.join(dir_path, filename)
            if os.path.isfile(file_path):
                mime_type, _ = mimetypes.guess_type(filename)
                documents.append({
                    'filename': filename,
                    'file_path': file_path,
                    'file_size': os.path.getsize(file_path),
                    'mime_type': mime_type or 'application/octet-stream',
                    'created_at': datetime.fromtimestamp(os.path.getctime(file_path)).isoformat(),
                    'modified_at': datetime.fromtimestamp(os.path.getmtime(file_path)).isoformat()
                })
        
        return documents

    def list_document_vector_chunks_paginated(self, document_uuid: str, page: int , per_page: int, keyword: str = '') -> List[Dict[str, str]]:
        """
        List document vector chunks with pagination.
        
        Args:
            document_uuid: UUID of the document
            page: Page number
            page_size: Number of chunks per page
            
        Returns:
            List of document vector chunk metadata dictionaries
        
        """

        print(document_uuid)

        document = Document.query.get(document_uuid)
        if not document:
            return []
        
        if keyword:
            document_vectors = DocumentVector.query.filter_by(document_uuid=document_uuid).filter(DocumentVector.text_content.contains(keyword)).all()
        else:
            document_vectors = DocumentVector.query.filter_by(document_uuid=document_uuid).all()
        
        """
        remove embediong attribute
        """
        for document_vector in document_vectors:
            document_vector.embedding = None

        total = len(document_vectors)
        total_pages = math.ceil(total / per_page)
        
        # Apply pagination
        start = (page - 1) * per_page
        end = start + per_page
        

        return document_vectors[start:end], document.to_dict(), total, total_pages

    def _extract_text(self, file_obj, content_type: str) -> str:
        """
        Extract text from a file object based on its content type.
        
        Args:
            file_obj: File-like object or file path
            content_type: MIME type of the file
            
        Returns:
            Extracted text content as a string
        """
        try:
            # Ensure we're at the start of the file
            if hasattr(file_obj, 'seek'):
                file_obj.seek(0)
                
            # PDF files
            if content_type == 'application/pdf' or (isinstance(content_type, str) and content_type.endswith('.pdf')):
                from PyPDF2 import PdfReader
                reader = PdfReader(file_obj)
                return "\n".join(page.extract_text() for page in reader.pages if page.extract_text())
            
            # Text-based files
            elif (isinstance(content_type, str) and 
                 (content_type.startswith('text/') or 
                  content_type in ['text/plain', 'text/markdown', 'text/x-markdown',
                                 'text/mdx', 'text/html', 'text/vtt'] or
                  any(ext in content_type.lower() for ext in ['.txt', '.md', '.mdx', '.html', '.htm', '.vtt', '.properties']))):
                content = file_obj.read()
                if isinstance(content, bytes):
                    return content.decode('utf-8', errors='replace')
                return content
                
            # CSV files
            elif (isinstance(content_type, str) and 
                 (content_type in ['text/csv', 'application/vnd.ms-excel'] or 
                  content_type.endswith('.csv'))):
                import csv
                import io
                content = file_obj.read()
                if isinstance(content, bytes):
                    content = content.decode('utf-8')
                text = io.StringIO()
                csv_reader = csv.reader(io.StringIO(content))
                for row in csv_reader:
                    text.write(' '.join(row) + '\n')
                return text.getvalue()
                
            # Excel files
            elif (isinstance(content_type, str) and 
                 (content_type in ['application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 
                                 'application/vnd.ms-excel'] or
                  any(ext in content_type.lower() for ext in ['.xlsx', '.xls']))):
                import pandas as pd
                return pd.read_excel(file_obj).to_string(index=False)
                
            # Word documents
            elif (isinstance(content_type, str) and 
                 (content_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document' or
                  content_type.endswith('.docx'))):
                from docx import Document
                doc = Document(file_obj)
                return "\n".join(paragraph.text for paragraph in doc.paragraphs if paragraph.text)
                
            # PowerPoint files
            elif (isinstance(content_type, str) and 
                 (content_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation' or
                  content_type.endswith('.pptx'))):
                from pptx import Presentation
                prs = Presentation(file_obj)
                text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text.append(shape.text)
                return "\n".join(text)

            else:
                # For unsupported types, try to read as text
                try:
                    content = file_obj.read()
                    if isinstance(content, bytes):
                        return content.decode('utf-8', errors='replace')
                    return content
                except:
                    return ""  # Return empty string for unsupported or unreadable files
                
        except Exception as e:
            # Log the error but don't fail the whole operation
            print(f"Error extracting text: {str(e)}")
            return ""  # Return empty string if extraction fails